import copy
import io
import math
import os
import re
from django.conf import settings
from django.http import HttpResponse
from django.utils import translation
from django.utils.translation import ugettext_lazy as _
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from democracy.enums import InitialSectionType
from democracy.models.section import SectionComment
from democracy.views.reports_v2.utils import (
    get_default_translation,
    get_formatted_hearing_timerange,
    get_powerpoint_title_font_size,
    get_selected_language,
)
from democracy.views.section_comment import SectionCommentSerializer

SLD_LAYOUT_MAIN_TITLE = 0
SLD_LAYOUT_SUBSECTION_TITLE = 1
SLD_LAYOUT_SECTION_COMMENTS = 2
SLD_LAYOUT_SECTION_POLL = 3
MAX_ROWS_PER_COMMENT_SLIDE = 11
MAX_CHARACTERS_PER_SINGLE_COMMENT_ROW = 75


class HearingReportPowerPoint:
    """Handles creating a powerpoint hearing report and serving it"""

    def __init__(self, json: object, context=None):
        self.json = json
        self.buffer = io.BytesIO()
        self.context = context
        self.initial_language = translation.get_language()
        self.used_language = get_selected_language(context['request'].query_params.get('lang'))
        translation.activate(self.used_language)
        self.prs = Presentation(self._get_theme_filename())

    def _get_theme_filename(self) -> str:
        dirname = os.path.dirname(__file__)
        theme = settings.HEARING_REPORT_THEME
        if theme == 'turku':
            return os.path.join(dirname, './powerpoint_templates/turku_reports_template.pptx')
        return os.path.join(dirname, './powerpoint_templates/whitelabel_reports_template.pptx')

    def _add_main_title_slide(self):
        title_slide_layout = self.prs.slide_layouts[SLD_LAYOUT_MAIN_TITLE]
        slide = self.prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        info = slide.placeholders[10]
        title_text = get_default_translation(self.json['title'], self.used_language)
        title.text = title_text
        title.text_frame.paragraphs[0].font.size = get_powerpoint_title_font_size(title_text)
        hearing_timerange = get_formatted_hearing_timerange(self.json['open_at'], self.json['close_at'])

        info.text = f'{_("Kerrokantasi hearing")} {hearing_timerange}' f'\n{_("Comments")} {self.json["n_comments"]}.'

    def _add_subsection_title_slide(self, section: dict):
        subsection_title_slide_layout = self.prs.slide_layouts[SLD_LAYOUT_SUBSECTION_TITLE]
        slide = self.prs.slides.add_slide(subsection_title_slide_layout)
        title = slide.shapes.title
        # sub section names use their title or type name if title doesnt exist
        section_title = get_default_translation(section['title'], self.used_language)
        if not section_title:
            section_title = section['type_name_singular']
        title.text = section_title
        title.text_frame.paragraphs[0].font.size = get_powerpoint_title_font_size(section_title, False)

    def _add_comment_slide(self, comments: list):
        comment_slide_layout = self.prs.slide_layouts[SLD_LAYOUT_SECTION_COMMENTS]
        slide = self.prs.slides.add_slide(comment_slide_layout)
        text_area = slide.placeholders[1].text_frame
        # text frame always contains one paragraph initially
        if len(comments) > 0:
            for index, comment in enumerate(comments):
                bullet = None
                if index == 0:
                    bullet = text_area.paragraphs[0]
                else:
                    bullet = text_area.add_paragraph()
                bullet.text = comment['content']
        else:
            bullet = text_area.paragraphs[0]
            bullet.text = f'{_("No comments")}'
            bullet.font.italic = True

    def _add_comment_slides(self, comments: dict):
        # calculate roughly how many text rows each comment will take
        comments_data = []
        for comment in comments:
            current_comment = comment
            # if comment is too long, truncate it
            comment_max_length = MAX_ROWS_PER_COMMENT_SLIDE * MAX_CHARACTERS_PER_SINGLE_COMMENT_ROW
            if len(current_comment['content']) > comment_max_length:
                current_comment['content'] = current_comment['content'][: comment_max_length - 3] + '...'

            comment_rows = int(math.ceil(len(current_comment['content']) / MAX_CHARACTERS_PER_SINGLE_COMMENT_ROW))
            comments_data.append({'comment': current_comment, 'comment_rows': comment_rows})

        # handle calculating how to insert comments into pages
        section_comment_pages = []
        comments_in_page = []
        used_row_counter = 0
        for comment_data in comments_data:
            if comment_data['comment_rows'] + used_row_counter <= MAX_ROWS_PER_COMMENT_SLIDE:
                # when comment fits in current page, insert it
                comments_in_page.append(comment_data['comment'])
                used_row_counter += comment_data['comment_rows']
            else:
                # when overflow would happen, start on a new page
                section_comment_pages.append(copy.deepcopy(comments_in_page))
                comments_in_page = []
                comments_in_page.append(comment_data['comment'])
                used_row_counter = comment_data['comment_rows']

        # add last filled page
        section_comment_pages.append(copy.deepcopy(comments_in_page))
        # create comment pages
        for page_comments in section_comment_pages:
            self._add_comment_slide(page_comments)

    def _add_poll_slide(self, poll: dict):
        poll_slide_layout = self.prs.slide_layouts[SLD_LAYOUT_SECTION_POLL]
        slide = self.prs.slides.add_slide(poll_slide_layout)

        chart_data = CategoryChartData()
        poll_options = poll['options']
        options_data = {'categories': [], 'n_answers': []}

        for option in poll_options:
            options_data['categories'].append(get_default_translation(option['text'], self.used_language))
            options_data['n_answers'].append(option['n_answers'])
        chart_data.categories = options_data['categories']
        chart_title = get_default_translation(poll['text'], self.used_language)
        chart_data.add_series(chart_title, options_data['n_answers'])

        x, y, cx, cy = Inches(1), Inches(0.5), Inches(8), Inches(4.5)
        slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)

        info = slide.placeholders[1]
        poll_type = poll['type']
        poll_total_answers = poll['n_answers']
        info.text = f'{_("Type")}: {_(poll_type)}, {_("total answers")}: {poll_total_answers}'

    def _add_section(self, section: dict):
        """
        Adds a hearing section with the following structure:
        - Title slide (main/sub)
        - Poll slides
        - Comment slides
        """
        # add title slide
        if section['type'] == 'main':
            self._add_main_title_slide()
        else:
            self._add_subsection_title_slide(section)

        # add poll slides
        for poll in section['questions']:
            self._add_poll_slide(poll)

        # add comment slides
        comments = [
            SectionCommentSerializer(c, context=self.context).data
            for c in SectionComment.objects.filter(section=section['id'])
        ]
        self._add_comment_slides(comments)

    def _get_pptx(self):
        sections = self.json['sections']
        for section in sections:
            if section['type'] != InitialSectionType.CLOSURE_INFO:
                self._add_section(section)

        self.prs.save(self.buffer)
        return self.buffer.getvalue()

    def get_response(self):
        """Returns http response with pptx content"""
        response = HttpResponse(
            self._get_pptx(), content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
        # remove special characters from filename to avoid potential file naming issues
        response['Content-Disposition'] = 'attachment; filename="{filename}.pptx"'.format(
            filename=re.sub(r"\W+|_", " ", get_default_translation(self.json['title'], self.used_language))
        )
        translation.activate(self.initial_language)
        return response
